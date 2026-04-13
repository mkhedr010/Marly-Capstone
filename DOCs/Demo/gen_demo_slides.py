"""
Generate ECG_CNN_Demo.pptx — 12-slide demo presentation.

Run from project root:
    pip install python-pptx
    python "DOCs/Demo/gen_demo_slides.py"

Output: DOCs/Demo/ECG_CNN_Demo.pptx

Figures embedded automatically:
  - figure_A_ecg_waveforms.png
  - figure_G_resource_utilization.png
  - figure_H_confusion_matrix.png
  - figure_I_model_comparison.png
  - figure_SYS_block_diagram.png
  - figure_QUANT_pipeline.png

Placeholder boxes left for:
  - [PHOTO] Python ECG visualizer screenshot  (slides 7, 10)
  - [PHOTO] Board LCD "Normal"                (slide 10)
  - [PHOTO] Board LCD "Abnormal"              (slide 10)
  - [EXPORT] Report Figure 4.1 — CNN arch     (slide 5)
  - [EXPORT] Report Figure 5.2 — UART FSM    (slide 8)
  - [EXPORT] Report Figure 5.4 — CNN FSM     (slide 9)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR   = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, '..', '..'))
FIG_DIR      = os.path.join(PROJECT_ROOT, 'DOCs', 'Poster', 'figures')
OUT_PATH     = os.path.join(SCRIPT_DIR, 'ECG_CNN_Demo.pptx')

def fig(name):
    return os.path.join(FIG_DIR, name)

# ── colors ─────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
BLUE  = RGBColor(0x15, 0x65, 0xC0)
RED   = RGBColor(0xDA, 0x25, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

# ── helpers ────────────────────────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text,
                font_size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
                bg_color=None, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if bg_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg_color
    return txb

def add_bullets(slide, left, top, width, height, items,
                font_size=20, color=None, bold_first=False):
    """items: list of strings. First char '•' added automatically."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f'  \u2022  {item}'
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and i == 0
        if color:
            run.font.color.rgb = color

def add_slide_header(slide, title, subtitle='', bg=NAVY):
    """Full-width navy header band with title text."""
    shp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    txb = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.6))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        txb2 = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.72), Inches(12.5), Inches(0.35))
        tf2  = txb2.text_frame
        p2   = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = LIGHT

def add_section_label(slide, left, top, width, label, color=BLUE):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.35))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    # underline via shape below
    line = slide.shapes.add_shape(1, Inches(left), Inches(top+0.3),
                                   Inches(width), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_placeholder(slide, left, top, width, height, label, color=BLUE):
    """Dashed placeholder box with centered label."""
    shp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shp.line.color.rgb = color
    shp.line.width = Pt(2)
    tf  = shp.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = color

def try_add_picture(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                  width=Inches(width), height=Inches(height))
        return True
    return False

def add_callout(slide, left, top, width, height, number, label, color=NAVY):
    shp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf  = shp.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = number
    r1.font.size = Pt(36)
    r1.font.bold = True
    r1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.color.rgb = LIGHT

# ── set up presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

# ECG waveform banner at bottom
if not try_add_picture(slide, fig('figure_A_ecg_waveforms.png'),
                        0, 5.4, 13.33, 2.1):
    add_placeholder(slide, 0, 5.4, 13.33, 2.1, '[Figure A: ECG waveforms banner]')

# title
add_textbox(slide, 0.5, 0.6, 12.3, 1.5,
            'Hardware Implementation of CNN\nas Part of SoC Design for ECG\nAnalysis and Classification',
            font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# author
add_textbox(slide, 0.5, 3.3, 12.3, 0.5,
            'Marly Barsoum',
            font_size=24, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 3.85, 12.3, 0.5,
            'Toronto Metropolitan University  |  Dept. of Electrical, Computer and Biomedical Engineering  |  2026',
            font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 4.4, 12.3, 0.4,
            'Capstone Design Project  —  FPGA-Based ECG Arrhythmia Classification',
            font_size=15, italic=True, color=RGBColor(0x64, 0xB5, 0xF6),
            align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATION
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Motivation & Problem', 'Why ECG classification on FPGA?')

add_bullets(slide, 0.5, 1.4, 7.5, 5.6, [
    'Cardiovascular disease is the #1 cause of death globally — 17.9 million deaths per year',
    'ECG is the primary non-invasive diagnostic tool — but manual reading requires specialists and is slow',
    'Automated classification enables real-time screening, wearable monitoring, and point-of-care diagnosis',
    'FPGAs offer deterministic low-latency inference at the edge — no cloud dependency, no CPU bottleneck',
    'Goal: Classify 8 arrhythmia types in 2.2 ms on a $100 FPGA board using a hardware CNN',
], font_size=21, color=DGRAY)

# 3 callout boxes right side
add_callout(slide, 8.4, 1.4, 2.2, 1.6, '17.9M', 'deaths/year\nfrom CVD', color=RED)
add_callout(slide, 10.8, 1.4, 2.2, 1.6, '360 Hz', 'ECG sampling\nrate', color=NAVY)
add_callout(slide, 8.4, 3.2, 2.2, 1.6, '2.2 ms', 'inference\nper beat', color=BLUE)
add_callout(slide, 10.8, 3.2, 2.2, 1.6, '95.75%', 'classification\naccuracy', color=GREEN)

add_textbox(slide, 8.4, 5.1, 4.6, 0.5,
            'MIT-BIH Arrhythmia Database: 8 classes, 360 Hz, 9,368 validation beats',
            font_size=13, italic=True, color=RGBColor(0x88, 0x88, 0x88))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'System Architecture', 'End-to-end data flow')

if not try_add_picture(slide, fig('figure_SYS_block_diagram.png'),
                        0.3, 1.25, 12.7, 3.9):
    add_placeholder(slide, 0.3, 1.25, 12.7, 3.9, '[Run gen_support_figures.py to generate this]')

add_bullets(slide, 0.5, 5.35, 12.3, 1.8, [
    'Python reads MIT-BIH ECG records  →  preprocesses to 12-bit  →  streams via RS-232 UART at 360 Hz',
    'FPGA receives samples  →  fills 128-sample buffer  →  runs ZolotyhNet CNN  →  drives LCD + LEDs',
    'PC visualizes the raw ECG waveform in real time on screen (matplotlib)',
], font_size=18, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASET & PREPROCESSING
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Dataset & Signal Preprocessing', 'MIT-BIH Arrhythmia Database')

# Left: bullets
add_section_label(slide, 0.4, 1.3, 6.0, 'MIT-BIH Database')
add_bullets(slide, 0.4, 1.75, 6.0, 2.5, [
    '48 records, 47 subjects, dual-lead ambulatory ECG',
    'Sampling rate: 360 Hz  |  MLII lead used',
    '8 arrhythmia classes: N, L, R, V, A, E, !, _',
    '128-sample window centered on R-peak (±64 samples ≈ 356 ms)',
    '90% training / 10% validation split  →  9,368 validation beats',
], font_size=18, color=DGRAY)

add_section_label(slide, 0.4, 4.3, 6.0, 'PC-Side Preprocessing')
add_bullets(slide, 0.4, 4.75, 6.0, 2.4, [
    'Raw ADC → physical mV:  (ADC − zero − baseline) / gain',
    'UART path: min-max normalize → 12-bit signed integer',
    'Two-byte framing: Byte1 = lower 8 bits, Byte2 = upper 4 bits',
], font_size=18, color=DGRAY)

# Right: Figure A
if not try_add_picture(slide, fig('figure_A_ecg_waveforms.png'),
                        6.6, 1.25, 6.5, 2.5):
    add_placeholder(slide, 6.6, 1.25, 6.5, 2.5, '[Figure A: ECG waveforms]')

add_textbox(slide, 6.6, 3.85, 6.5, 0.45,
            'Figure 1: 128-sample beat windows — Normal (N), PVC (V), LBBB (L) — MIT-BIH @ 360 Hz',
            font_size=12, italic=True, color=RGBColor(0x77, 0x77, 0x77))

# Class table
add_textbox(slide, 6.6, 4.5, 6.5, 0.35,
            'The 8 arrhythmia classes:', font_size=15, bold=True, color=NAVY)

classes = [
    ('N', 'Normal sinus beat'), ('L', 'Left Bundle Branch Block'),
    ('R', 'Right Bundle Branch Block'), ('V', 'PVC'),
    ('A', 'Atrial Premature Beat'), ('E', 'Ventricular Escape Beat'),
    ('!', 'Ventricular Flutter'), ('_', 'Unclassified'),
]
row_h = 0.26
for i, (code, name) in enumerate(classes):
    y = 4.95 + i * row_h
    col = NAVY if i % 2 == 0 else BLUE
    add_textbox(slide, 6.6, y, 0.5, row_h, code,
                font_size=14, bold=True, color=col)
    add_textbox(slide, 7.2, y, 5.9, row_h, name,
                font_size=13, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CNN ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'ZolotyhNet — Dual-Path 1D CNN',
                 'Novel architecture designed for FPGA resource constraints')

add_bullets(slide, 0.4, 1.3, 5.8, 4.0, [
    'Upper path (Conv1D): 5 layers + 4 MaxPool — extracts local beat morphology',
    'Lower path (FC): 3 Linear layers — captures global signal statistics',
    'Fusion: element-wise ADD of both paths (8 values each)',
    'Classifier: Linear(8→8) → Argmax → 3-bit class output (0–7)',
    '~14,700 parameters — 34× smaller than EcgResNet34 (~500K)',
    'Quantized to Q8.8 fixed-point — fits entirely in on-chip M4K ROM',
], font_size=19, color=DGRAY)

# Key stats boxes
for i, (num, lbl) in enumerate([('14.7K', 'parameters'), ('5+3', 'layer depth'),
                                  ('8', 'output classes'), ('Q8.8', 'quantization')]):
    add_callout(slide, 0.4 + i*1.45, 5.6, 1.3, 1.55, num, lbl, color=NAVY)

# Right: CNN architecture placeholder (user inserts from report)
add_placeholder(slide, 6.5, 1.25, 6.5, 5.95,
                '[INSERT: CNN Architecture Diagram\n\nExport Figure 4.1 from Final_Report_Draft.md\n(ZolotyhNet dual-path block diagram)\n\nUse Mermaid Live or VS Code Preview\nto export as PNG, then insert here]',
                color=BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — QUANTIZATION PIPELINE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Fixed-Point Quantization',
                 'From PyTorch float32 weights to FPGA hardware')

if not try_add_picture(slide, fig('figure_QUANT_pipeline.png'),
                        0.3, 1.2, 12.7, 3.4):
    add_placeholder(slide, 0.3, 1.2, 12.7, 3.4, '[Run gen_support_figures.py to generate this]')

add_section_label(slide, 0.5, 4.85, 5.8, 'Q8.8 Format')
add_bullets(slide, 0.5, 5.3, 5.8, 1.9, [
    '16-bit signed: 1 sign + 7 integer + 8 fractional bits',
    'Range: −128.0 to +127.996  |  Resolution: 1/256 ≈ 0.004',
    'Conversion: round(weight × 256).astype(int16)',
], font_size=17, color=DGRAY)

add_section_label(slide, 6.8, 4.85, 6.0, 'MAC Arithmetic')
add_bullets(slide, 6.8, 5.3, 6.0, 1.9, [
    'Q8.8 × Q8.8 → Q16.16 product (32-bit accumulator)',
    'After accumulation: right-shift 8 bits → back to Q8.8',
    '18 .mif files → 18 weight_rom instances in VHDL',
], font_size=17, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PYTHON INFRASTRUCTURE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Python Infrastructure',
                 'Two-thread streaming + real-time visualization')

# Left content
add_section_label(slide, 0.4, 1.3, 5.8, 'ecg_stream_visualize.py')
add_bullets(slide, 0.4, 1.75, 5.8, 3.0, [
    'ECGStreamer: daemon thread — reads ECG file, streams at 360 Hz (time.sleep(1/360))',
    'ECGVisualizer: main thread — matplotlib animation (required by Tkinter)',
    'Thread-safe queue (maxsize=2000) between threads',
    'Display path: raw mV values → scrolling ECG plot on PC',
    'UART path: min-max normalize → 12-bit signed → two-byte frame → serial port',
], font_size=18, color=DGRAY)

add_section_label(slide, 0.4, 4.95, 5.8, 'Two-Byte UART Framing')
add_bullets(slide, 0.4, 5.4, 5.8, 1.8, [
    'Byte 1 = lower 8 bits of 12-bit sample',
    'Byte 2 = upper 4 bits, zero-padded to 8 bits',
    '115,200 baud  |  8N1  |  CLKS_PER_BIT = 434',
], font_size=17, color=DGRAY)

# Right: Python visualizer screenshot placeholder
add_placeholder(slide, 6.5, 1.25, 6.5, 5.95,
                '[YOUR SCREENSHOT HERE]\n\nRun:\npython python/ecg_stream_visualize.py\n\nCapture a screenshot while the ECG\nis scrolling on screen\n(Full matplotlib window, no other apps visible)',
                color=GREEN)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FPGA: UART & INPUT BUFFER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'FPGA Hardware: UART & Input Buffer',
                 'uart_receiver.vhd  +  buffer_128.vhd')

add_section_label(slide, 0.4, 1.3, 6.0, 'UART Receiver — Two Nested FSMs')
add_bullets(slide, 0.4, 1.75, 6.0, 2.4, [
    'FSM 1 (byte-level): IDLE → START_BIT → DATA_BITS (×8) → STOP_BIT',
    'FSM 2 (sample assembly): WAIT_BYTE1 → WAIT_BYTE2 → pulse sample_valid',
    '115,200 baud  |  8N1  |  CLKS_PER_BIT = 434  |  50 MHz clock',
    'uart_error asserted on bad stop bit',
], font_size=18, color=DGRAY)

add_section_label(slide, 0.4, 4.3, 6.0, 'Input Buffer — buffer_128.vhd')
add_bullets(slide, 0.4, 4.75, 6.0, 2.4, [
    '128-entry circular dual-port M4K RAM',
    'Write: 12-bit raw ADC  |  Read: 16-bit Q8.8',
    '3-clock pipeline: sign-extend 12→16 → arithmetic right-shift 4',
    'Triggers CNN when 128 samples accumulated AND CNN is IDLE',
], font_size=18, color=DGRAY)

# Right: UART FSM diagram placeholder
add_placeholder(slide, 6.5, 1.25, 6.5, 5.95,
                '[INSERT: UART FSM Diagram]\n\nExport Figure 5.2 from Final_Report_Draft.md\n(UART Receiver State Machines — both FSMs)\n\nMermaid → PNG via mermaid.live\nthen insert here',
                color=BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FPGA: CNN ACCELERATOR
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'FPGA Hardware: CNN Accelerator',
                 'zolotyhnet_top  +  conv1d_engine  +  linear_engine')

add_section_label(slide, 0.4, 1.3, 6.0, '13-State Top-Level FSM')
add_bullets(slide, 0.4, 1.75, 6.0, 1.6, [
    'IDLE → CONV1→2→3→4→5 → LINEAR1→2→3 → FUSION → CLASSIFIER → ARGMAX → OUTPUT',
    'OUTPUT holds result_valid=\'1\' for 1,000 cycles (~20 µs)',
    'FUSION: upper[i] + lower[i] element-wise (not concatenation)',
], font_size=17, color=DGRAY)

add_section_label(slide, 0.4, 3.5, 6.0, 'Time-Multiplexed MAC Engines')
add_bullets(slide, 0.4, 3.95, 6.0, 2.1, [
    'Conv1D engine: 9-state FSM — handles all 5 conv layers sequentially',
    'Linear engine: 10-state FSM — handles Linear1, Linear2, Linear3, Classifier',
    'Both use 32-bit accumulator (Q16.16) → right-shift 8 → Q8.8',
    'ReLU inline in WRITE_OUTPUT (conv) / APPLY_RELU (linear) states',
], font_size=17, color=DGRAY)

add_section_label(slide, 0.4, 6.15, 6.0, '18 Weight ROMs')
add_bullets(slide, 0.4, 6.6, 6.0, 0.7, [
    'weight_rom.vhd — ramstyle="M4K"  |  ram_init_file → .mif  |  61% M4K used',
], font_size=17, color=DGRAY)

# Right: CNN FSM placeholder
add_placeholder(slide, 6.5, 1.25, 6.5, 5.95,
                '[INSERT: CNN FSM + Engine Diagrams]\n\nExport Figure 5.4 from Final_Report_Draft.md\n(ZolotyhNet VHDL Top-Level Hardware Architecture)\n\nAND/OR Figure 3.10 (Conv1D Engine FSM)\n\nMermaid → PNG via mermaid.live\nthen insert here',
                color=RED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LIVE DEMO
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

add_textbox(slide, 0.5, 0.15, 12.3, 0.8,
            'LIVE DEMONSTRATION',
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 0.85, 12.3, 0.45,
            'Streaming MIT-BIH ECG records from PC → FPGA → LCD classification in real time',
            font_size=17, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# Board photos — Normal (left) and Abnormal (right)
add_placeholder(slide, 0.3, 1.45, 4.7, 3.8,
                '[BOARD PHOTO]\nLCD: "Normal"\n\nHigh-res photo of DE2 board\nwith Normal ECG streaming',
                color=GREEN)
add_placeholder(slide, 5.2, 1.45, 4.7, 3.8,
                '[BOARD PHOTO]\nLCD: "Abnormal"\n\nHigh-res photo of DE2 board\nwith Abnormal ECG streaming',
                color=RED)

# Python visualizer screenshot
add_placeholder(slide, 10.1, 1.45, 2.9, 3.8,
                '[SCREENSHOT]\nPython ECG\nvisualizer\nrunning',
                color=BLUE)

# Labels
for x, lbl, col in [(2.65, '"NORMAL"', GREEN), (7.55, '"ABNORMAL"', RED)]:
    add_textbox(slide, x-1.5, 5.4, 3.0, 0.5, lbl,
                font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)

# Key numbers at bottom
add_callout(slide, 0.3,  6.05, 2.8, 1.2, '2.2 ms', 'CNN inference', color=BLUE)
add_callout(slide, 3.3,  6.05, 2.8, 1.2, '360 Hz', 'stream rate',   color=BLUE)
add_callout(slide, 6.3,  6.05, 2.8, 1.2, '357 ms', 'end-to-end',    color=BLUE)
add_callout(slide, 9.3,  6.05, 3.7, 1.2, '~32 records tested', 'MIT-BIH validation', color=RGBColor(0x2E,0x7D,0x32))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Classification Results',
                 'ZolotyhNet on 9,368-sample MIT-BIH validation set')

# Key accuracy number — big
add_textbox(slide, 0.3, 1.2, 6.5, 1.1,
            '95.75%',
            font_size=60, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.3, 2.2, 6.5, 0.45,
            'Overall Accuracy  |  8,970 / 9,368 beats correct',
            font_size=16, color=DGRAY, align=PP_ALIGN.CENTER)

add_bullets(slide, 0.3, 2.85, 6.5, 2.5, [
    'Normal (N): 99.37%  — dominant class, 70% of dataset',
    'LBBB (L): 96.41%  |  RBBB (R): 94.76%  — well learned',
    'PVC (V): 83.50%  |  APB (A): 56.63%  — partial confusion',
    'VEB (E): 0.00%  |  Flutter (!): 4.26%  — class imbalance',
    '2.2 ms inference  |  Linear1 (128→64) is the dominant layer at ~45% of cycles',
], font_size=18, color=DGRAY)

# Confusion matrix
if not try_add_picture(slide, fig('figure_H_confusion_matrix.png'),
                        0.3, 5.35, 5.5, 1.95):
    add_placeholder(slide, 0.3, 5.35, 5.5, 1.95, '[Figure H: Confusion Matrix]')

# Model comparison
if not try_add_picture(slide, fig('figure_I_model_comparison.png'),
                        6.9, 1.2, 6.1, 4.0):
    add_placeholder(slide, 6.9, 1.2, 6.1, 4.0, '[Figure I: Model Comparison]')

add_textbox(slide, 6.9, 5.25, 6.1, 0.6,
            'Trade-off: 3.6% accuracy reduction vs. best software model → enables full FPGA deployment',
            font_size=14, italic=True, color=DGRAY)

# Resource chart
if not try_add_picture(slide, fig('figure_G_resource_utilization.png'),
                        6.9, 5.9, 6.1, 1.35):
    add_placeholder(slide, 6.9, 5.9, 6.1, 1.35, '[Figure G: Resource Utilization]')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY CONTRIBUTIONS & CONCLUSIONS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_slide_header(slide, 'Key Contributions & Conclusions', '')

add_section_label(slide, 0.4, 1.2, 7.5, 'What We Built')
add_bullets(slide, 0.4, 1.65, 7.5, 4.5, [
    '\u2713  ZolotyhNet: a novel 14,700-parameter dual-path CNN — designed for FPGA resource constraints',
    '\u2713  Q8.8 fixed-point pipeline: PyTorch training → extract_weights.py → 18 MIF files → on-chip ROM',
    '\u2713  Time-multiplexed VHDL hardware: 1 Conv1D + 1 Linear MAC engine → 13% LEs, 61% M4K',
    '\u2713  95.75% accuracy on 8 arrhythmia classes — 2.2 ms inference per 128-sample beat window',
    '\u2713  Complete Python infrastructure: MIT-BIH reader, preprocessing, streaming, visualization',
    '\u2713  End-to-end validated: ~32 ECG records streamed to board — LCD and LED outputs correct',
], font_size=20, color=DGRAY)

# Resource chart left-bottom
if not try_add_picture(slide, fig('figure_G_resource_utilization.png'),
                        0.4, 6.1, 5.5, 1.15):
    add_placeholder(slide, 0.4, 6.1, 5.5, 1.15, '[Figure G]')

# Key numbers right
add_section_label(slide, 8.2, 1.2, 4.8, 'System Highlights')
for i, (num, lbl) in enumerate([
    ('13%',    'Logic Elements used'),
    ('61%',    'M4K RAM blocks used'),
    ('14.7K',  'CNN parameters'),
    ('2.2 ms', 'per-beat inference'),
    ('95.75%', 'classification accuracy'),
]):
    y = 1.7 + i * 0.9
    col = RED if i == 4 else NAVY
    add_callout(slide, 8.2, y, 4.8, 0.75, num, lbl, color=col)

# ── save ───────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"\nDone! Saved: {OUT_PATH}")
print("\nPlaceholders remaining (open PowerPoint and insert):")
print("  Slide  5: CNN architecture diagram       (Report Fig 4.1)")
print("  Slide  7: Python visualizer screenshot   (your photo)")
print("  Slide  8: UART FSM diagram               (Report Fig 5.2)")
print("  Slide  9: CNN FSM/engine diagrams        (Report Fig 5.4)")
print("  Slide 10: Board photo 'Normal'           (your photo)")
print("  Slide 10: Board photo 'Abnormal'         (your photo)")
print("  Slide 10: Python visualizer screenshot   (your photo)")
